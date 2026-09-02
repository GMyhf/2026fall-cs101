#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""verify_courseware.py —— 《计算概论（B）》讲义与课件的一致性闸门。

用法:
    python3 tools/verify_courseware.py            # 第 1–7 项，约 10 秒
    python3 tools/verify_courseware.py --render   # 加第 8 项渲染检查（需 libreoffice + pdftotext）

检查项:
    1  配对        每周恰好一份 .md + 同名 .pptx + content/wNN.py，W01–W16 无缺漏
    2  元数据      讲义有一级标题、Updated 时间戳、Compiled by、仓库 URL
    3  课程安排    讲义与课件声明的「主题与学习重点」与课程指南表格逐字一致
    4  链接        所有本地 .md / .pptx / .py 相对链接可达
    5  语法        讲义里所有 ```python 代码块 + courseware/*.py 能被 ast.parse
    6  可重生成    课件能从 content/ 重新生成；页数与 README 一致，
                   且重建产物与已提交的 .pptx **逐段文本相同**（防止源改了没重建）
    7  题号题名    讲义引用的 OJ 题号↔题名与仓库内既有语料一致（离线）
    8  渲染        逐页检查文字未越出版心 + 中文字体已嵌入（--render）
    9  机考规格    三份样卷同一规格：6 题 / 112 分钟；
                   题目标题、难度梯度和讲义正文不写固定分值分配
   10  版面标记    放映稿的非等宽文字里不得印出 `**` / 反引号；
                   bullet 不得以空白（含全角空格）开头
   11  列对齐      等宽图里不得有跨过不同中文字数的对齐列
                   （汉字 1 em vs Consolas 约 0.55 em，补空格对不出来）

只用标准库 + python-pptx（第 6、10 项）。退出码 0 表示全绿。
"""

import argparse
import ast
import os
import re
import subprocess
import sys
import tempfile
import unicodedata
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
COURSEWARE = ROOT / 'courseware'
CONTENT = COURSEWARE / 'content'
GUIDE = ROOT / 'Introduction_to_Computing_B_Course_Guide.md'

WEEKS = [f'{i:02d}' for i in range(1, 17)]

# 第 15、16 周共用课程安排表里的「第 15-16 周」一行
GUIDE_ROW_FOR_WEEK = {w: w for w in WEEKS}
GUIDE_ROW_FOR_WEEK['15'] = '15-16'
GUIDE_ROW_FOR_WEEK['16'] = '15-16'

# 有意偏离既有题名语料的引用：题号 -> {讲义里的写法: 理由}
# 形状是「题号 + 具体别名 + 理由」，不做整号豁免 —— 整号豁免会让该号的任何错叫法蒙混过关。
TITLE_WHITELIST = {
    # 例：'02734': {'十进制转八进制': '平台标题的另一种常见写法'},
}

# ---------------------------------------------------------------------------
# 已联网核实过的平台题名：题号 -> (权威题名, 出处)
#
# ⚠️ 这张表的存在理由，是本项目踩过的一个坑：
#     **语料本身会过时。** 第 7 项拿 `2025fall-cs101/` 当语料，
#     而 2025 年的材料里把 12559 写成「最大最小整数 v0.3」——
#     平台上的实际题名是「最大最小整数」。语料错了，闸门就会**祝福错误答案**，
#     再怎么收紧匹配规则也没用（变异自检实测：改回 v0.3 时收紧后的规则依然全绿）。
#
# 所以：**凡在此表登记的题号，以本表为准，语料被忽略。**
# 表里的每一条都必须来自真正的联网核实，并注明出处（任务号 / 提交号）。
VERIFIED_TITLES = {
    '12559': ('最大最小整数',
              'T-008 联网核实（Codex，b2e6e3e）：平台标题无 "v0.3" 后缀，'
              '2025 语料的写法已过时'),
}

failures = []
notes = []


def fail(check, msg):
    failures.append(f'[{check}] {msg}')


def note(msg):
    notes.append(msg)


def read(path):
    return path.read_text(encoding='utf-8')


# ---------------------------------------------------------------- 公共解析
def week_files():
    """返回 {周次: (md 路径, pptx 路径, content 路径)}；缺失的用 None。"""
    # build_all 在导入时会 import deck（需要 python-pptx）；这里只读常量，直接解析源码
    src = read(COURSEWARE / 'build_all.py')
    tree = ast.parse(src)
    mapping = None
    for node in ast.walk(tree):
        if isinstance(node, ast.Assign):
            for t in node.targets:
                if isinstance(t, ast.Name) and t.id == 'WEEKS':
                    mapping = ast.literal_eval(node.value)
    if mapping is None:
        fail('1 配对', 'build_all.py 里找不到 WEEKS 映射')
        return {}
    out = {}
    for wk in WEEKS:
        stem = mapping.get(wk)
        if stem is None:
            fail('1 配对', f'build_all.WEEKS 缺少第 {wk} 周')
            out[wk] = (None, None, None)
            continue
        md = COURSEWARE / f'{stem}.md'
        pptx = COURSEWARE / f'{stem}.pptx'
        py = CONTENT / f'w{wk}.py'
        out[wk] = (md if md.exists() else None,
                   pptx if pptx.exists() else None,
                   py if py.exists() else None)
    return out


def guide_rows():
    """解析课程指南的「课程安排」表：{'01': '主题与学习重点原文', ...}"""
    if not GUIDE.exists():
        fail('3 课程安排', f'找不到课程指南 {GUIDE.name}')
        return {}
    rows = {}
    for line in read(GUIDE).splitlines():
        m = re.match(r'^\|\s*第\s*([0-9]+(?:\s*-\s*[0-9]+)?)\s*周\s*\|(.*?)\|\s*$', line)
        if not m:
            continue
        key = m.group(1).replace(' ', '')
        topic = m.group(2).strip()
        if '-' in key:
            lo, hi = key.split('-')
            rows[f'{int(lo)}-{int(hi)}'] = topic
        else:
            rows[f'{int(key):02d}'] = topic
    return rows


PY_BLOCK = re.compile(r'```python\n(.*?)```', re.S)
CODE_ANY = re.compile(r'```.*?```', re.S)
INLINE_CODE = re.compile(r'`[^`\n]*`')


# ---------------------------------------------------------------- 检查 1
def check_pairing(files):
    for wk in WEEKS:
        md, pptx, py = files.get(wk, (None, None, None))
        if md is None:
            fail('1 配对', f'第 {wk} 周缺少讲义 .md')
        if pptx is None:
            fail('1 配对', f'第 {wk} 周缺少课件 .pptx')
        if py is None:
            fail('1 配对', f'第 {wk} 周缺少内容模块 content/w{wk}.py')
    # 反向：courseware 下不应有孤儿 md / pptx
    known_md = {f.name for f, _, _ in files.values() if f}
    known_pptx = {f.name for _, f, _ in files.values() if f}
    for f in COURSEWARE.glob('*.md'):
        if f.name != 'README.md' and f.name not in known_md:
            fail('1 配对', f'孤儿讲义（不在 build_all.WEEKS 中）：{f.name}')
    for f in COURSEWARE.glob('*.pptx'):
        # Microsoft Office creates transient lock files such as
        # `~$202610_ADS_W08_Recursion.pptx` while a deck is open.
        # They are neither courseware nor Git assets.
        if f.name.startswith('~$'):
            continue
        if f.name not in known_pptx:
            fail('1 配对', f'孤儿课件：{f.name}')
    for f in CONTENT.glob('w*.py'):
        if f.stem[1:] not in WEEKS:
            fail('1 配对', f'孤儿内容模块：{f.name}')


# ---------------------------------------------------------------- 检查 2
def check_metadata(files):
    for wk in WEEKS:
        md = files.get(wk, (None,))[0]
        if md is None:
            continue
        text = read(md)
        head = '\n'.join(text.splitlines()[:8])
        if not re.match(r'^#\s+第\s*\d+\s*周\s', text):
            fail('2 元数据', f'{md.name}: 首行不是「# 第N周 ...」形式的一级标题')
        if not re.search(r'\*Updated\s+\d{4}-\d{2}-\d{2}[^*]*GMT\+8\*', head):
            fail('2 元数据', f'{md.name}: 缺少 *Updated YYYY-MM-DD GMT+8* 时间戳')
        if '*Compiled by' not in head:
            fail('2 元数据', f'{md.name}: 缺少 *Compiled by ...*')
        if 'github.com/GMyhf/2026fall-cs101' not in head:
            fail('2 元数据', f'{md.name}: 头部缺少仓库 URL')


# ---------------------------------------------------------------- 检查 3
def content_meta(py_path):
    """从 content/wNN.py 中静态取出 META 字典。"""
    tree = ast.parse(read(py_path))
    for node in tree.body:
        if isinstance(node, ast.Assign):
            for t in node.targets:
                if isinstance(t, ast.Name) and t.id == 'META':
                    return ast.literal_eval(node.value)
    return None


def check_syllabus(files):
    rows = guide_rows()
    if not rows:
        return
    for wk in WEEKS:
        key = GUIDE_ROW_FOR_WEEK[wk]
        if key not in rows:
            fail('3 课程安排', f'课程指南表格里找不到第 {key} 周这一行')
            continue
        expected = rows[key]

        md, _, py = files.get(wk, (None, None, None))
        if md is not None:
            text = read(md)
            m = re.search(r'^>\s*\*\*主题与学习重点\*\*：(.*)$', text, re.M)
            if not m:
                fail('3 课程安排', f'{md.name}: 缺少 > **主题与学习重点**： 一行')
            elif m.group(1).strip() != expected:
                fail('3 课程安排',
                     f'{md.name}: 主题与学习重点与课程指南不一致\n'
                     f'        讲义: {m.group(1).strip()}\n'
                     f'        指南: {expected}')
        if py is not None:
            meta = content_meta(py)
            if meta is None:
                fail('3 课程安排', f'{py.name}: 找不到 META')
                continue
            info = meta.get('info') or []
            decl = [s for s in info if s.startswith('主题与学习重点：')]
            if not decl:
                fail('3 课程安排', f'{py.name}: META["info"] 缺少「主题与学习重点：」一项')
            elif decl[0][len('主题与学习重点：'):].strip() != expected:
                fail('3 课程安排',
                     f'{py.name}: 课件 META 与课程指南不一致\n'
                     f'        课件: {decl[0][len("主题与学习重点："):].strip()}\n'
                     f'        指南: {expected}')


# ---------------------------------------------------------------- 检查 4
LINK = re.compile(r'\[[^\]]*\]\(([^)\s]+)\)')


def check_links(files):
    for wk in WEEKS:
        md = files.get(wk, (None,))[0]
        if md is None:
            continue
        text = CODE_ANY.sub('', read(md))          # 代码块里的括号不是链接
        for target in LINK.findall(text):
            if target.startswith(('http://', 'https://', 'mailto:', '#')):
                continue
            target = target.split('#')[0]
            if not target:
                continue
            if not target.endswith(('.md', '.pptx', '.py')):
                continue
            resolved = (md.parent / target).resolve()
            if not resolved.exists():
                fail('4 链接', f'{md.name}: 死链 {target}')
    for f in [COURSEWARE / 'README.md', ROOT / 'README.md']:
        if not f.exists():
            continue
        text = CODE_ANY.sub('', read(f))
        for target in LINK.findall(text):
            if target.startswith(('http://', 'https://', 'mailto:', '#')):
                continue
            target = target.split('#')[0]
            if not target or not target.endswith(('.md', '.pptx', '.py')):
                continue
            if not (f.parent / target).resolve().exists():
                fail('4 链接', f'{f.name}: 死链 {target}')


# ---------------------------------------------------------------- 检查 5
def check_syntax(files):
    total = 0
    for wk in WEEKS:
        md = files.get(wk, (None,))[0]
        if md is None:
            continue
        for i, block in enumerate(PY_BLOCK.findall(read(md)), start=1):
            total += 1
            try:
                ast.parse(block)
            except SyntaxError as e:
                fail('5 语法', f'{md.name}: 第 {i} 个 python 代码块语法错误：{e}')
    for py in sorted(list(COURSEWARE.glob('*.py')) + list(CONTENT.glob('*.py'))):
        try:
            ast.parse(read(py))
        except SyntaxError as e:
            fail('5 语法', f'{py.relative_to(ROOT)}: 语法错误：{e}')
    note(f'共校验 {total} 个讲义 python 代码块')


# ---------------------------------------------------------------- 检查 6
README_ROW = re.compile(
    r'^\|\s*(\d+)\s*\|\s*`([^`]+)`\s*\|\s*(\d+)\s*\|', re.M)


def pptx_texts(path):
    """抽出 .pptx 里全部可见文字（含表格单元格），用于比对两份课件是否同一内容。

    不比对 XML 或字节：zip 时间戳、元素顺序每次生成都会变，只有可见文本是稳定的。
    """
    from pptx import Presentation
    out = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if getattr(shape, 'has_table', False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        out.append(cell.text)
    return out


def check_regenerate(files):
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        note('⚠️ 未安装 python-pptx，**跳过**第 6 项（可重生成）—— '
             '课件与 content/ 的一致性本次未验')
        return
    sys.path.insert(0, str(COURSEWARE))
    sys.path.insert(0, str(CONTENT))
    import importlib
    import deck

    declared = {}
    readme = COURSEWARE / 'README.md'
    if readme.exists():
        for wk, stem, pages in README_ROW.findall(read(readme)):
            declared[f'{int(wk):02d}'] = (stem, int(pages))
    else:
        fail('6 可重生成', 'courseware/README.md 不存在')

    with tempfile.TemporaryDirectory() as tmp:
        for wk in WEEKS:
            py = files.get(wk, (None, None, None))[2]
            if py is None:
                continue
            mod = importlib.import_module(f'w{wk}')
            importlib.reload(mod)
            out = Path(tmp) / f'w{wk}.pptx'
            pages = deck.build(mod.META, mod.SLIDES, str(out))
            if wk not in declared:
                fail('6 可重生成', f'README 文件清单里没有第 {wk} 周')
                continue
            stem, want = declared[wk]
            got_stem = files[wk][1].stem if files[wk][1] else '(缺)'
            if stem != got_stem:
                fail('6 可重生成',
                     f'第 {wk} 周 README 文件名 {stem} 与实际 {got_stem} 不一致')
            if pages != want:
                fail('6 可重生成',
                     f'第 {wk} 周页数不一致：重新生成 {pages} 页，README 声明 {want} 页')

            # ⚠️ 只比页数是不够的：改了 content/wNN.py 却忘了重建课件时，
            #    页数往往一模一样，仓库里的 .pptx 就这样悄悄过期。
            #    必须把重新生成的产物与已提交的课件**逐段文本**比对。
            committed = files[wk][1]
            if committed is None:
                continue
            fresh_texts, old_texts = pptx_texts(out), pptx_texts(committed)
            if fresh_texts != old_texts:
                diff = next((f'第 {i + 1} 段：仓库「{o[:40]}」≠ 重建「{n[:40]}」'
                             for i, (o, n) in enumerate(zip(old_texts, fresh_texts))
                             if o != n),
                            f'段数不同：仓库 {len(old_texts)} vs 重建 {len(fresh_texts)}')
                fail('6 可重生成',
                     f'第 {wk} 周课件与 content/w{wk}.py 不一致 —— '
                     f'源改过但 .pptx 没重建？\n        {diff}\n'
                     f'        修法：cd courseware && python3 build_all.py {wk}')


# ---------------------------------------------------------------- 检查 7
CORPUS_GLOBS = ['2025fall-cs101/*.md', 'ADS_problem_list_at_*.md',
                'ADS_matrices.md']
# 讲义里的引用形态：**E02750: 鸡兔同笼**、`02760: 数字三角形`、表格里的 02773 等
# 5 位题号 + 冒号 + 题名。
#   - 前面不能是数字或小数点：排除 3.14159 这类误匹配；
#   - 允许前面是中文或字母：语料里有"练习01742: Coins"、"OJ02806:公共子序列"、
#     "E02750: 鸡兔同笼"等多种写法，\b 在 CJK 与数字之间并不成立；
#   - 后面不能再跟数字：排除长数字串的尾部；
#   - 题名在 `]`（Markdown 链接起点）处截断，但**允许括号** ——
#     早先把 `(`/`（` 也排除掉，导致语料里「汉诺塔问题(Tower of Hanoi)」
#     被截成「汉诺塔问题」，反过来把讲义里正确的全称判成不符。
REF = re.compile(
    r'(?<![0-9.])(\d{5})(?!\d)\s*[:：]\s*([^\n,，。|`*<\]]{1,30})')


# 表格形态的引用：`| ... | 题名 | 编号 | ... |`（题名在前、题号在后，中间没有冒号）。
# W16 的备选题库、各周的作业表都是这个形状 —— REF 那条正则**完全看不到它们**，
# 而它们恰恰是本仓库题号引用最密集的地方（实测 71 处）。
TABLE_ID = re.compile(r'^[ETM]?(\d{5})$')
TITLE_HEADERS = ('题目', '题名')


def table_refs(text):
    r"""从 Markdown 表格里抽 (题号, 题名)。

    **按表头列名定位**，不用"题号紧邻前一格"这种启发式 ——
    本仓库里有 `| 形状 | 排序键 | 例题 |`、`| 形态 | 递归时传 | 典型题 |` 这类表，
    题号前面一格是排序键 / 递归参数，不是题名（实测会造成 5 处误报）。

    只处理**表头里含「题目」或「题名」列**的表：题名取该列，题号取同一行里
    形如 `[ETM]?\d{5}` 的单元格。合并行（如 `LC 20 / 02694`）不匹配，自动跳过。
    """
    out = []
    header_idx = None
    for line in text.splitlines():
        line = line.strip()
        if not line.startswith('|'):
            header_idx = None           # 表格结束
            continue
        cells = [c.strip().strip('`').strip() for c in line.strip('|').split('|')]
        if set(''.join(cells)) <= set('-: '):
            continue                    # 分隔行 |---|---|
        if header_idx is None:
            for k, c in enumerate(cells):
                if any(h in c for h in TITLE_HEADERS):
                    header_idx = k
                    break
            else:
                header_idx = -1         # 这张表没有题名列，整表跳过
            continue
        if header_idx < 0 or header_idx >= len(cells):
            continue
        title = cells[header_idx].strip('*` ').strip()
        if not title or title.isdigit():
            continue
        for c in cells:
            m = TABLE_ID.match(c)
            if m:
                out.append((m.group(1), title))
                break
    return out


def normalize_title(s):
    """归一化题名，使不同书写方式可以精确比较。

    语料是从既有 Markdown 里正则抽出来的，会带上链接残留（如 `最大最小整数]`），
    必须在这里洗干净 —— 否则就只能靠"子串容差"兜，而子串容差会**放过真正的题名漂移**
    （T-008 实测：过时的 `最大最小整数 v0.3` 与正确题名互为子串，被静默通过）。
    """
    s = s.strip()
    s = re.sub(r'[\s　]+', '', s)
    s = s.replace('（', '(').replace('）', ')')
    s = s.replace('，', ',').replace('“', '"').replace('”', '"')
    s = s.strip('])】》')                 # Markdown 链接残留
    s = s.rstrip('.,:;、')
    return s.lower()


def build_title_corpus():
    """从仓库既有材料里抽取 {题号: {已知题名}}。"""
    corpus = {}
    for pattern in CORPUS_GLOBS:
        for path in sorted(ROOT.glob(pattern)):
            for num, title in REF.findall(read(path)):
                t = normalize_title(title)
                if not t or t.isdigit():
                    continue
                corpus.setdefault(num, set()).add(t)
    return corpus


def check_problem_titles(files):
    corpus = build_title_corpus()
    if not corpus:
        fail('7 题号题名', '未能从仓库既有材料中抽取到任何题号↔题名语料')
        return
    checked = unknown = 0
    verified = sum(
        1 for wk in WEEKS if files.get(wk, (None,))[0]
        for num, _t in (REF.findall(read(files[wk][0]))
                        + table_refs(read(files[wk][0])))
        if num in VERIFIED_TITLES)
    for wk in WEEKS:
        md = files.get(wk, (None,))[0]
        if md is None:
            continue
        refs = REF.findall(read(md)) + table_refs(read(md))
        for num, title in refs:
            t = normalize_title(title)
            if not t or t.isdigit():
                continue
            if num in VERIFIED_TITLES:
                # 已联网核实过：以权威题名为准，忽略（可能过时的）语料
                checked += 1
                official, source = VERIFIED_TITLES[num]
                if t == normalize_title(official):
                    continue
                if title.strip() in TITLE_WHITELIST.get(num, {}):
                    continue
                fail('7 题号题名',
                     f'{md.name}: 题号 {num} 的题名「{title.strip()}」'
                     f'与已联网核实的平台题名「{official}」不符\n'
                     f'        出处：{source}')
                continue
            known = corpus.get(num)
            if not known:
                unknown += 1
                continue                       # 语料里没有该题号，无从判定
            checked += 1
            if t in known:
                continue
            # ⚠️ 这里**不做子串容差**。曾经写过 `if t in k or k in t: continue`，
            #    结果过时题名 `最大最小整数 v0.3` 与正确题名互为子串，漂移被静默放过
            #    （由 T-008 联网核实才发现）。归一化洗掉链接残留后，
            #    精确匹配对当前 73 处引用零误报。
            allowed = TITLE_WHITELIST.get(num, {})
            if title.strip() in allowed:
                continue
            fail('7 题号题名',
                 f'{md.name}: 题号 {num} 的题名「{title.strip()}」'
                 f'与仓库既有语料不符，已知：{sorted(known)}')
    note(f'题号↔题名：逐处比对 {checked} 处'
         f'（其中 {verified} 处以联网核实的权威题名为准）；'
         f'{unknown} 处题号不在既有语料中（无从判定）')


# ---------------------------------------------------------------- 检查 8
CJK_FONT_FAMILIES = (
    'simsun', 'simhei', 'simkai', 'simfang', 'nsimsun',
    'microsoftyahei', 'msyh', 'microsoftjhenghei',
    'kaiti', 'fangsong', 'youyuan', 'lisu', 'stsong', 'stkaiti', 'stheiti',
    'stfangsong', 'stxihei', 'stzhongsong', 'songti', 'heiti', 'pingfang',
    'hiraginosansgb', 'notosanssc', 'notoserifsc', 'notosanstc',
    'notoseriftc', 'notosanscjk', 'notoserifcjk', 'sourcehansans',
    'sourcehanserif', 'wenquanyi', 'wqy', 'droidsansfallback',
    'arplumingcn', 'arplukaicn', 'ipagothic', 'ipamincho', 'nanumgothic',
    'malgun', 'meiryo', 'msgothic', 'msmincho', 'yugothic', 'yumincho',
)


def has_cjk_font(pdffonts_output):
    """纯包含式判定：只认已知的 CJK 字体家族名，不用通用词兜底。"""
    for line in pdffonts_output.splitlines()[2:]:
        parts = line.split()
        if not parts:
            continue
        name = re.sub(r'[^a-z]', '', parts[0].lower())
        if any(fam in name for fam in CJK_FONT_FAMILIES):
            return True
    return False


# deck.py 的版面几何（pt，16:9 = 960x540）：正文区与页脚的分界
BODY_TOP_PT = 1.42 * 72
BODY_BOTTOM_PT = BODY_TOP_PT + 5.32 * 72      # = 485.3
# 页脚**文本**的实测顶端（LibreOffice 渲染下 yMin≈497.5，略高于文本框的 498.2）：
# 阈值必须按实测取，不能按理论值 —— 差 0.7pt 就会把每一页的页脚都误判成溢出。
FOOTER_TEXT_TOP_PT = 496.0
BODY_TOL_PT = 4.0                             # 字形下伸部等的容差


def intrudes_into_footer(y_min_pt, y_max_pt):
    """正文文字是否越出版心底部并侵入页脚区（坐标已换算到 540pt 版面）。

    抽成具名函数是为了能**脱离 LibreOffice 直接回归**这条判据 ——
    它的阈值来自实测（页脚文本 yMin≈497.5，而文本框理论值是 498.2），
    差 0.7pt 就会把每一页的页脚都误判成溢出。
    """
    return (y_min_pt < FOOTER_TEXT_TOP_PT
            and y_max_pt > BODY_BOTTOM_PT + BODY_TOL_PT)


def check_render(files):
    soffice = None
    for cand in ('soffice', 'libreoffice'):
        try:
            subprocess.run([cand, '--version'], capture_output=True, check=True)
            soffice = cand
            break
        except (FileNotFoundError, subprocess.CalledProcessError):
            continue
    if soffice is None:
        fail('8 渲染', '找不到 libreoffice / soffice，无法执行渲染检查')
        return
    try:
        subprocess.run(['pdftotext', '-v'], capture_output=True)
    except FileNotFoundError:
        fail('8 渲染', '找不到 pdftotext（poppler-utils）')
        return

    pptxs = [files[wk][1] for wk in WEEKS if files[wk][1]]
    with tempfile.TemporaryDirectory() as tmp:
        proc = subprocess.run(
            [soffice, '--headless', '--convert-to', 'pdf', '--outdir', tmp]
            + [str(p) for p in pptxs],
            capture_output=True, text=True, timeout=1800)
        pdfs = sorted(Path(tmp).glob('*.pdf'))
        if len(pdfs) != len(pptxs):
            fail('8 渲染',
                 f'转换未全部成功：期望 {len(pptxs)} 份 PDF，实得 {len(pdfs)} 份；'
                 f'soffice 退出码 {proc.returncode}')
            return
        total_pages = 0
        for pdf in pdfs:
            fonts = subprocess.run(['pdffonts', str(pdf)],
                                   capture_output=True, text=True).stdout
            if not has_cjk_font(fonts):
                fail('8 渲染', f'{pdf.name}: PDF 未嵌入任何已知中文字体，'
                               f'渲染结果可能是方框（本机字体环境限制）')
            # 逐页文本框位置检查
            bbox = subprocess.run(['pdftotext', '-bbox', str(pdf), '-'],
                                  capture_output=True, text=True).stdout
            pages = re.findall(
                r'<page width="([\d.]+)" height="([\d.]+)">(.*?)</page>',
                bbox, re.S)
            total_pages += len(pages)
            for pno, (pw, ph, body) in enumerate(pages, start=1):
                pw, ph = float(pw), float(ph)
                # pdftotext 的坐标按 PDF 点数给出，换算到 deck.py 的 720x540 版面
                sy = 540.0 / ph
                for x, y, x2, y2 in re.findall(
                        r'<word xMin="([\d.]+)" yMin="([\d.]+)" '
                        r'xMax="([\d.]+)" yMax="([\d.]+)"', body):
                    x, y, x2, y2 = float(x), float(y), float(x2), float(y2)
                    if x2 > pw + 0.5 or y2 > ph + 0.5 or x < -0.5 or y < -0.5:
                        fail('8 渲染',
                             f'{pdf.name} 第 {pno} 页：文字越出页面 '
                             f'({x},{y})-({x2},{y2}) vs {pw}x{ph}')
                        break
                    # ⚠️ 只比对页面边界是不够的：正文压到页脚上仍然在页内。
                    #    W16 曾有 3 页代码压住页脚而这里报"0 处越界"，
                    #    直到把判据改成**版心**底部才暴露。
                    if intrudes_into_footer(y * sy, y2 * sy):
                        fail('8 渲染',
                             f'{pdf.name} 第 {pno} 页：正文越出版心底部并侵入页脚区'
                             f'（文字底 {y2 * sy:.1f}pt > 版心底 {BODY_BOTTOM_PT:.1f}pt）')
                        break
        note(f'渲染检查：{len(pdfs)} 份 PDF，共 {total_pages} 页')


# ---------------------------------------------------------------- 检查 9
# 三次月考与期末上机考试是**同一规格**：6 题 / 112 分钟。这个规格由任课教师给定
# （不在课程指南表里）；固定分值分配另有核算办法，课件与讲义均不得自行写死。
# 以前只有人眼盯着：W05 写 5 题、W14 写 4 题、W16 同时写着
# 「2025 秋为 112 分钟」和「约 120 分钟」，全程无人报警。
EXAM_MINUTES = 112
EXAM_QUESTION_COUNT = 6
EXAM_PAPERS = {'05': '10 月月考样卷', '14': '12 月月考样卷', '16': '期末上机考试样卷'}

# 被 112 分钟取代的旧写法。「8 小时」是课外学习时间，不能误伤，所以只卡 2 小时。
STALE_DURATION = re.compile(r'约?\s*120\s*分钟|(?<![0-9])2\s*小时')
EXAM_HEAD = re.compile(r'^##\s*T(\d+)\.\s*(.+?)\s*$', re.M)
EXAM_HEAD_SCORE = re.compile(r'^##\s*T\d+\..*?（\d+\s*分）\s*$', re.M)
SCORE_ALLOCATION = re.compile(r'分值分配')

# 难度梯度与题目标题只回答"多难"和"考什么"，不回答"几分"。
# 梯度表里再列一列 15分/20分，就成了同一事实的第五个写法；而且机考之后
# **另有成绩核算办法**，这一列会跟真正的核算口径打架。所以梯度表里禁止出现分数。
LADDER_MD = re.compile(r'\*\*难度梯度\*\*：\s*\n*(?:```\n(.*?)```|([^\n]*))', re.S)
LADDER_SCORE = re.compile(r'\d+\s*分(?!钟)')


def _strings(obj):
    """把 slide 元组里嵌套的所有字符串摊平（表格是 list of list）。"""
    if isinstance(obj, str):
        yield obj
    elif isinstance(obj, (list, tuple)):
        for x in obj:
            yield from _strings(x)


def exam_ladders(md_text, py_path):
    """逐行产出讲义与课件里的难度梯度表：(是讲义吗, 行)。

    ⚠️ 课件侧**不能只认 `ascii` 那一种**：T-017 把梯度表改成了 `table`，
    判据若还盯着 ascii，表格单元格里再写一列分值就会静默放过。
    这里改成"标题带'难度梯度'的整张页，所有字符串都算梯度表内容"。
    """
    for m in LADDER_MD.finditer(md_text):
        for ln in (m.group(1) or m.group(2) or '').splitlines():
            yield True, ln
    for slide in deck_slides(py_path) if py_path else ():
        is_ladder = len(slide) >= 2 and '难度梯度' in str(slide[1])
        for text in _strings(slide):
            if is_ladder or '难度梯度' in text:
                for ln in text.splitlines():
                    yield False, ln


def check_exam_spec(files):
    for wk, (md, _pptx, py) in sorted(files.items()):
        for path in (md, py):
            if path is None:
                continue
            for m in STALE_DURATION.finditer(read(path)):
                fail('9 机考规格',
                     f'{path.name}: 出现过时的考试时长写法 {m.group(0)!r}；'
                     f'机考一律 {EXAM_MINUTES} 分钟')

    for wk, label in EXAM_PAPERS.items():
        md, _pptx, py = files.get(wk, (None, None, None))
        if md is None:
            fail('9 机考规格', f'第 {wk} 周讲义缺失，无法校验{label}')
            continue
        text = read(md)

        if f'{EXAM_MINUTES} 分钟' not in text:
            fail('9 机考规格',
                 f'{md.name}: {label}没有写明「{EXAM_MINUTES} 分钟」')

        heads = EXAM_HEAD.findall(text)
        got_ids = [int(i) for i, _t in heads]
        if got_ids != list(range(1, EXAM_QUESTION_COUNT + 1)):
            fail('9 机考规格',
                 f'{md.name}: {label}的题号应为 T1..T{EXAM_QUESTION_COUNT}，'
                 f'实际 {got_ids}')

        scored_heads = EXAM_HEAD_SCORE.findall(text)
        if scored_heads:
            fail('9 机考规格',
                 f'{md.name}: {label}的题目标题不应写分数；'
                 '固定分值另有核算办法')

        if SCORE_ALLOCATION.search(text):
            fail('9 机考规格',
                 f'{md.name}: {label}不应出现「分值分配」；固定分值另有核算办法')
        if py is not None and f'{EXAM_MINUTES} 分钟' not in read(py):
            fail('9 机考规格',
                 f'{py.name}: 课件没有写明「{EXAM_MINUTES} 分钟」')

        for is_md, ln in exam_ladders(text, py):
            if LADDER_SCORE.search(ln):
                fail('9 机考规格',
                     f'{(md if is_md else py).name}: {label}的难度梯度里出现分数'
                     f'{ln.strip()!r}；固定分值另有核算办法，'
                     f'机考成绩另有核算办法')

    note(f'机考规格：3 份样卷均为 {EXAM_QUESTION_COUNT} 题 / {EXAM_MINUTES} 分钟')


# ---------------------------------------------------------------- 检查 10
# 版面标记：`**` 与反引号是**源里的记号**，绝不该出现在放映稿上。
#
# 起因（T-013 复核）：`deck.py` 的 `_SEGMENT` 把 `**...**` 整段吞掉、不再往里拆，
# 于是 ``**负数不能用 `bin(x)`**`` 里的反引号被原样印了出来。
# 全仓库 8 处这种写法，在 T-007（439 页）、T-011（8 页）、T-013（29 页）
# **三轮逐页人眼复核里全部漏过** —— 一个小反引号，人眼就是不敏感。
# 同轮还发现一条用全角空格伪造的"续行"，被渲成了一个空 bullet。
# 两类都属于"看得见但看不出"的缺陷，只能由代码守。
MONO_FONT_NAME = 'Consolas'
LEAK = re.compile(r'\*\*|`')
LEADING_BLANK = re.compile(r'^[\s\u3000]')


def deck_slides(py_path):
    """静态取出 content/wNN.py 的 SLIDES（不 import，避免副作用）。"""
    tree = ast.parse(read(py_path))
    for node in tree.body:
        if isinstance(node, ast.Assign):
            for t in node.targets:
                if isinstance(t, ast.Name) and t.id == 'SLIDES':
                    try:
                        return ast.literal_eval(node.value)
                    except ValueError:
                        return None
    return None


def check_markup(files):
    # --- A 层：源里的 bullet / 表格单元格不得以空白（含全角空格）开头。
    #     用全角空格伪造"续行"会渲成一个只有项目符号、没有内容的空 bullet。
    for wk, (_md, _pptx, py) in sorted(files.items()):
        if py is None:
            continue
        slides = deck_slides(py)
        if slides is None:
            fail('10 版面标记', f'{py.name}: 无法静态解析 SLIDES')
            continue
        for i, sl in enumerate(slides):
            kind = sl[0] if sl else '?'
            if kind == 'bullets':
                items = sl[2] if len(sl) > 2 else []
                for it in items:
                    if isinstance(it, str) and LEADING_BLANK.search(it):
                        fail('10 版面标记',
                             f'{py.name} 第 {i} 张「{sl[1]}」：bullet 以空白开头 '
                             f'{it[:24]!r} —— 伪造的续行会渲成空项目符号，'
                             f'请拆成两条完整 bullet')
            elif kind == 'table':
                for row in (sl[2] if len(sl) > 2 else []):
                    for cell in row:
                        if isinstance(cell, str) and LEADING_BLANK.search(cell):
                            fail('10 版面标记',
                                 f'{py.name} 第 {i} 张「{sl[1]}」：表格单元格以空白开头 '
                                 f'{cell[:24]!r}')

    # --- B 层：产物里的**非等宽** run 不得含 `**` 或反引号。
    #     检查对象必须是渲染结果 —— 源写法千变万化，泄漏只有一种表现。
    try:
        from pptx import Presentation
    except ImportError:
        note('⚠️ 未安装 python-pptx，第 10 项**只做了源侧 A 层**；'
             '产物侧的标记泄漏本次未验')
        return
    leaked = 0
    for wk, (_md, pptx, _py) in sorted(files.items()):
        if pptx is None:
            continue
        for n, slide in enumerate(Presentation(str(pptx)).slides, 1):
            for shape in slide.shapes:
                frames = []
                if shape.has_text_frame:
                    frames.append(shape.text_frame)
                if getattr(shape, 'has_table', False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            frames.append(cell.text_frame)
                for tf in frames:
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.name == MONO_FONT_NAME:
                                continue          # 代码页原样输出，** 是幂运算符
                            if LEAK.search(run.text):
                                leaked += 1
                                fail('10 版面标记',
                                     f'{pptx.name} 第 {n} 页：非等宽文字里印出了 '
                                     f'Markdown 记号 {run.text[:40]!r}')
    if not leaked:
        note('版面标记：16 份课件的非等宽文字里没有 `**` / 反引号泄漏')


# ---------------------------------------------------------------- 检查 11
# 等宽版面的列对齐：**汉字宽度对不上任何等宽拉丁字体**，
# 所以"靠补空格把一列对齐"这件事，只要那一列前面隔着中文就一定做不到。
#
# 算一遍就清楚：汉字的字身宽是 1 em，Consolas 的字身宽约 0.55 em。
# 两个空格 = 1.10 em ≠ 一个汉字 = 1.00 em。每多一个汉字就少 0.10 em，
# 各行中文字数不同，同一"显示列"落到纸上就是不同的位置。
#
# 起因（T-016）：W05 p10 难度梯度 6 行的 `--`，
# 在 Microsoft PowerPoint 16.112.3 下**完全不在同一列**（人眼当场判不通过）；
# 本机 LibreOffice 复测同一页，6 个 `--` 的 x 极差 **36.85 pt**。
# 同一页的 W16 p13 只有 2 行、前缀汉字数只差 2，极差 1.02 pt，
# 于是**看上去是对的** —— 它不是对的，只是漂得还看不出来。
# 判据因此不看"漂了多少"，只看"这一列是不是跨过了不同的中文字数"。
#
# 判据里不含"漂多少 pt"，是有意的：pt 要靠 LibreOffice 渲染才量得到，
# 而字体替换会让它失真；而"前缀汉字数不同"是**源里就能判定的确定事实**。
# `=` / `+` / `#` 是 T-024 补进来的：写 W10 那张三角形图时我自己踩了 —— 
# 「6-7-0-4 = 17 奇」把一个汉字夹在两列中间，右边那一列当场歪掉，
# 而当时的记号表里没有 `=`，闸门一声没吭。补完之后全仓库零新增（`+` 那一处
# 落在 w01「四个概念层层包含」，本来就在 KNOWN_RAGGED 上）。
ALIGN_TOKENS = ('--', '——', '->', '→', '=>', '⇒', '<-', '←', '|', '=', '+', '#')

# 已知未修的 7 张图（都是"方框里写中文、右边框对不齐"这一种），见 T-018。
# 逐页实测的 `|` 列 x 极差：W01 p17 9.2pt、W02 p5 7.8pt、W03 p6 2.4pt、
# W03 p10 3.5pt、W08 p10 10.0pt、W13 p11 13.4pt、W15 p12 0.1pt。
# 都远小于判不通过的那 36.85pt，但性质完全一样。挂账在这里而不是删掉判据，
# 是为了**新写的 ascii 图立刻变红**，同时不假装这 7 张已经修好。
# 名单按（模块, 标题）记而不按页序号记：插一张新页就会把序号全推后，
# 那样挂账会**悄悄挪到别的图上**，而真正欠着的那张反而变红。
KNOWN_RAGGED = {
    ('w01', '四个概念层层包含'),
    ('w02', '虚拟机的层次'),
    ('w03', '图灵机模型（1936）'),
    ('w03', '五大部件'),
    ('w08', '进程的虚拟地址空间'),
    ('w13', '虚拟地址空间'),
    ('w15', 'XOR：为什么需要"深度"'),
}


def _display_width(text):
    """显示宽度：中日韩全宽字算 2 格，其余算 1 格（源里排版就是按这个数的）。"""
    return sum(2 if unicodedata.east_asian_width(c) in 'WF' else 1 for c in text)


def _wide_count(text):
    return sum(1 for c in text if unicodedata.east_asian_width(c) in 'WF')


def check_alignment(files):
    checked = ragged = 0
    for wk, (_md, _pptx, py) in sorted(files.items()):
        if py is None:
            continue
        slides = deck_slides(py)
        if slides is None:
            continue
        for si, slide in enumerate(slides):
            if slide[0] != 'ascii' or len(slide) < 3:
                continue
            checked += 1
            cols = {}
            for line in str(slide[2]).splitlines():
                for tok in ALIGN_TOKENS:
                    for m in re.finditer(re.escape(tok), line):
                        if m.start() == 0:
                            continue          # 行首，前缀是空的，天然对齐
                        prefix = line[:m.start()]
                        cols.setdefault((tok, _display_width(prefix)), set()).add(
                            _wide_count(prefix))
            worst = max((max(v) - min(v) for v in cols.values() if len(v) > 1),
                        default=0)
            if not worst:
                continue
            if (py.stem, str(slide[1])) in KNOWN_RAGGED:
                ragged += 1
                continue
            fail('11 列对齐',
                 f'{py.name} 第 {si + 2} 页「{slide[1]}」：等宽图里有一列'
                 f'跨过了不同的中文字数（最多相差 {worst} 个汉字），'
                 f'真实渲染必然对不齐。改用表格，或把要对齐的那一列'
                 f'挪到中文**前面**（前缀只剩 ASCII 就能对准）')
    note(f'列对齐：{checked} 张等宽图，'
         f'{ragged} 张为 T-018 挂账的历史遗留，其余无跨中文的对齐列')


# ---------------------------------------------------------------- 主流程
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--render', action='store_true',
                    help='附加第 8 项渲染检查（需 libreoffice + pdftotext）')
    opts = ap.parse_args()
    if os.environ.get('VERIFY_RENDER') == '1':
        opts.render = True

    files = week_files()
    check_pairing(files)
    check_metadata(files)
    check_syllabus(files)
    check_links(files)
    check_syntax(files)
    check_regenerate(files)
    check_problem_titles(files)
    if opts.render:
        check_render(files)
    check_exam_spec(files)
    check_markup(files)
    check_alignment(files)

    print('=' * 68)
    for n in notes:
        print('  ·', n)
    if failures:
        print(f'\n✗ 闸门未通过：{len(failures)} 项')
        for f in failures:
            print('  -', f)
        return 1
    print('\n✓ 闸门全部通过'
          + ('（含渲染检查）' if opts.render else '（未含渲染检查，加 --render 开启）'))
    return 0


if __name__ == '__main__':
    sys.exit(main())
